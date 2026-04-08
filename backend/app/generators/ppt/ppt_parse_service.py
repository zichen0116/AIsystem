"""
PPT生成模块 - 文件解析服务
用于解析 PPT/PDF 文件（翻新功能）
"""
import os
import re
import logging
from typing import Optional, Tuple

logger = logging.getLogger(__name__)


class PptParseService:
    """PPT/PDF 文件解析服务"""

    def __init__(self, mineru_token: str = None):
        """
        初始化解析服务

        Args:
            mineru_token: MinerU API token（用于解析 PDF/PPTX）
        """
        self.mineru_token = mineru_token or os.getenv("MINERU_TOKEN", "")
        self.mineru_api_base = os.getenv("MINERU_API_BASE", "https://mineru.net")

    async def parse_file(self, file_path: str, filename: str) -> Tuple[Optional[str], Optional[str]]:
        """
        解析文件，返回 markdown 内容和错误信息

        Args:
            file_path: 文件路径
            filename: 文件名

        Returns:
            (markdown_content, error_message)
        """
        ext = os.path.splitext(filename)[1].lower()

        if ext in ['.pdf']:
            return await self._parse_pdf(file_path)
        elif ext in ['.pptx', '.ppt']:
            return await self._parse_pptx(file_path)
        elif ext in ['.txt', '.md']:
            return await self._parse_text(file_path)
        elif ext in ['.xlsx', '.csv']:
            return await self._parse_spreadsheet(file_path)
        else:
            return None, f"不支持的文件类型: {ext}"

    async def _parse_pdf(self, file_path: str) -> Tuple[Optional[str], Optional[str]]:
        """解析 PDF 文件"""
        if not self.mineru_token:
            return None, "MinerU token 未配置，请联系管理员"

        try:
            # 导入必要的库
            import requests

            # 步骤1: 获取上传 URL
            upload_url_api = f"{self.mineru_api_base}/api/v4/file-urls/batch"
            headers = {"Authorization": f"Bearer {self.mineru_token}"}
            payload = {"files": [{"file_name": os.path.basename(file_path), "file_size": os.path.getsize(file_path)}]}

            resp = requests.post(upload_url_api, headers=headers, json=payload, timeout=30)
            if resp.status_code != 200:
                return None, f"MinerU 上传失败: {resp.text}"

            upload_info = resp.json()["data"][0]
            upload_url = upload_info["upload_url"]
            file_key = upload_info["file_key"]

            # 步骤2: 上传文件
            with open(file_path, "rb") as f:
                file_data = f.read()
            upload_resp = requests.put(upload_url, data=file_data, timeout=60)
            if upload_resp.status_code not in (200, 201):
                return None, f"MinerU 文件上传失败: {upload_resp.status_code}"

            # 步骤3: 获取解析结果
            result_url_template = f"{self.mineru_api_base}/api/v4/extract-results/batch/{file_key}"

            # 轮询获取结果
            import asyncio
            for _ in range(30):  # 最多等待 30 * 10 = 300 秒
                await asyncio.sleep(10)
                result_resp = requests.get(result_url_template, headers=headers, timeout=30)
                if result_resp.status_code == 200:
                    result_data = result_resp.json()
                    if result_data.get("status") == "completed":
                        # 下载 markdown 结果
                        markdown_url = result_data["data"]["result_files"][0]["url"]
                        md_resp = requests.get(markdown_url, timeout=60)
                        if md_resp.status_code == 200:
                            return md_resp.text, None
                        else:
                            return None, f"下载 markdown 失败: {md_resp.status_code}"
                    elif result_data.get("status") == "failed":
                        return None, f"MinerU 解析失败: {result_data.get('message', '未知错误')}"

            return None, "解析超时，请稍后重试"

        except Exception as e:
            logger.exception("PDF 解析异常")
            return None, f"PDF 解析异常: {str(e)}"

    async def _parse_pdf_v4(self, file_path: str) -> Tuple[Optional[str], Optional[str]]:
        """MinerU v4 PDF parsing flow."""
        try:
            import asyncio
            import io
            import requests
            import zipfile

            upload_url_api = f"{self.mineru_api_base}/api/v4/file-urls/batch"
            headers = {"Authorization": f"Bearer {self.mineru_token}"}
            payload = {
                "files": [{
                    "name": os.path.basename(file_path),
                    "is_ocr": True,
                    "data_id": f"pdf-{int(os.path.getmtime(file_path))}-{os.path.getsize(file_path)}",
                }]
            }

            resp = requests.post(upload_url_api, headers=headers, json=payload, timeout=30)
            if resp.status_code != 200:
                return None, f"MinerU upload failed: {resp.text}"

            upload_payload = resp.json()
            if upload_payload.get("code") != 0:
                return None, f"MinerU upload failed: {upload_payload.get('msg', upload_payload)}"

            data = upload_payload.get("data") or {}
            batch_id = data.get("batch_id")
            file_urls = data.get("file_urls") or []
            if not batch_id or not file_urls:
                return None, f"MinerU upload response missing batch_id/file_urls: {upload_payload}"

            with open(file_path, "rb") as f:
                file_data = f.read()
            upload_resp = requests.put(file_urls[0], data=file_data, timeout=60)
            if upload_resp.status_code not in (200, 201):
                return None, f"MinerU file upload failed: {upload_resp.status_code}"

            result_url = f"{self.mineru_api_base}/api/v4/extract-results/batch/{batch_id}"
            for _ in range(30):
                await asyncio.sleep(10)
                result_resp = requests.get(result_url, headers=headers, timeout=30)
                if result_resp.status_code != 200:
                    continue

                result_data = result_resp.json()
                if result_data.get("code") != 0:
                    return None, f"MinerU parse failed: {result_data.get('msg', result_data)}"

                extract_result = ((result_data.get("data") or {}).get("extract_result") or [])
                current = extract_result[0] if isinstance(extract_result, list) and extract_result else extract_result
                if not isinstance(current, dict):
                    continue

                state = str(current.get("state") or "").lower()
                if state in {"done", "success", "completed"}:
                    markdown_url = current.get("full_md_url")
                    if markdown_url:
                        md_resp = requests.get(markdown_url, timeout=60)
                        if md_resp.status_code == 200:
                            return md_resp.text, None
                        return None, f"Download markdown failed: {md_resp.status_code}"

                    zip_url = current.get("full_zip_url")
                    if zip_url:
                        zip_resp = requests.get(zip_url, timeout=120)
                        if zip_resp.status_code != 200:
                            return None, f"Download zip failed: {zip_resp.status_code}"
                        with zipfile.ZipFile(io.BytesIO(zip_resp.content)) as zf:
                            candidates = [
                                name for name in zf.namelist()
                                if name.lower().endswith("full.md")
                            ]
                            if not candidates:
                                candidates = [name for name in zf.namelist() if name.lower().endswith(".md")]
                            if not candidates:
                                return None, "MinerU zip does not contain markdown output"
                            with zf.open(candidates[0]) as md_file:
                                return md_file.read().decode("utf-8", errors="ignore"), None

                    return None, "MinerU result missing markdown/zip url"

                if state in {"failed", "error"}:
                    return None, f"MinerU parse failed: {current.get('err_msg') or current.get('state')}"

            return None, "PDF parse timed out, please retry later"

        except Exception as e:
            logger.exception("PDF v4 parse exception")
            return None, f"PDF parse exception: {str(e)}"

    async def _parse_pptx(self, file_path: str) -> Tuple[Optional[str], Optional[str]]:
        """解析 PPTX 文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation(file_path)
            markdown_lines = ["# PPT 内容\n"]

            for i, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"\n## 第 {i} 页\n")

                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 简单的清理
                        text = re.sub(r'\s+', ' ', text)
                        markdown_lines.append(f"- {text}\n")

                # 如果没有提取到文本
                if len(markdown_lines) == 0 or markdown_lines[-1] == f"\n## 第 {i} 页\n":
                    markdown_lines.append("- （此页无文本内容）\n")

            return "".join(markdown_lines), None

        except Exception as e:
            logger.exception("PPTX 解析异常")
            return None, f"PPTX 解析异常: {str(e)}"

    async def _parse_text(self, file_path: str) -> Tuple[Optional[str], Optional[str]]:
        """解析文本文件"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            return content, None
        except Exception as e:
            return None, f"文本文件读取失败: {str(e)}"

    async def _parse_spreadsheet(self, file_path: str) -> Tuple[Optional[str], Optional[str]]:
        """解析电子表格文件"""
        try:
            import pandas as pd

            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            # 转换为 markdown 表格
            markdown = df.to_markdown(index=False)
            return markdown, None

        except Exception as e:
            return None, f"电子表格解析失败: {str(e)}"

    async def split_pdf_to_pages(self, pdf_path: str, output_dir: str) -> list[str]:
        """
        将 PDF 拆分为单独页面

        Args:
            pdf_path: PDF 文件路径
            output_dir: 输出目录

        Returns:
            拆分后的 PDF 页面路径列表
        """
        output_paths = []

        try:
            import fitz  # PyMuPDF
        except ModuleNotFoundError:
            from PyPDF2 import PdfReader, PdfWriter

            with open(pdf_path, "rb") as f:
                reader = PdfReader(f)
                for page_num, page in enumerate(reader.pages, start=1):
                    writer = PdfWriter()
                    writer.add_page(page)
                    output_path = os.path.join(output_dir, f"page_{page_num}.pdf")
                    with open(output_path, "wb") as out:
                        writer.write(out)
                    output_paths.append(output_path)
            return output_paths

        pdf_document = fitz.open(pdf_path)

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            # 创建新的 PDF 文档
            new_doc = fitz.open()
            new_doc.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)

            output_path = os.path.join(output_dir, f"page_{page_num + 1}.pdf")
            new_doc.save(output_path)
            new_doc.close()
            output_paths.append(output_path)

        pdf_document.close()
        return output_paths


# 单例
_parse_service: Optional[PptParseService] = None


def get_ppt_parse_service() -> PptParseService:
    """获取解析服务单例"""
    global _parse_service
    if _parse_service is None:
        _parse_service = PptParseService()
    return _parse_service
