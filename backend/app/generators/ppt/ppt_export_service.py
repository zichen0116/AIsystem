"""
PPT export service.

Provides the existing image-based exporters and keeps editable PPTX export
inside the current repository by overlaying editable text on top of slide
images.
"""
import io
import logging
import os
from typing import Optional


logger = logging.getLogger(__name__)


def get_page_size_inches(aspect_ratio: str) -> tuple:
    sizes = {
        "16:9": (13.333, 7.5),
        "4:3": (10, 7.5),
        "1:1": (7.5, 7.5),
    }
    return sizes.get(aspect_ratio, (13.333, 7.5))


class PptExportService:
    """PPT export service."""

    @staticmethod
    def create_editable_pptx(
        image_paths: list[str],
        output_file: str = None,
        aspect_ratio: str = "16:9",
        pages_data: list[dict] = None,
    ) -> tuple[Optional[bytes], dict]:
        pptx_bytes = PptExportService.create_pptx_from_images(
            image_paths,
            output_file=output_file,
            aspect_ratio=aspect_ratio,
            pages_data=pages_data,
            add_text_layer=True,
        )
        return pptx_bytes, {
            "mode": "beta",
            "warnings": [],
            "warning_details": {},
        }

    @staticmethod
    def create_pptx_from_images(
        image_paths: list[str],
        output_file: str = None,
        aspect_ratio: str = "16:9",
        pages_data: list[dict] = None,
        add_text_layer: bool = False,
    ) -> Optional[bytes]:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        page_w, page_h = get_page_size_inches(aspect_ratio)
        prs.slide_width = Inches(page_w)
        prs.slide_height = Inches(page_h)

        for idx, image_path in enumerate(image_paths):
            if not os.path.exists(image_path):
                logger.warning("Image does not exist: %s", image_path)
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

            if add_text_layer and pages_data and idx < len(pages_data):
                page = pages_data[idx]
                title_text = page.get("title") or ""
                desc_text = page.get("description") or ""
                notes_text = page.get("notes") or ""

                if title_text:
                    title_box = slide.shapes.add_textbox(
                        Inches(page_w * 0.05),
                        Inches(page_h * 0.04),
                        Inches(page_w * 0.9),
                        Inches(page_h * 0.12),
                    )
                    title_frame = title_box.text_frame
                    title_frame.word_wrap = True
                    title_paragraph = title_frame.paragraphs[0]
                    title_paragraph.alignment = PP_ALIGN.LEFT
                    title_run = title_paragraph.add_run()
                    title_run.text = title_text
                    title_run.font.size = Pt(28)
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                if desc_text:
                    desc_box = slide.shapes.add_textbox(
                        Inches(page_w * 0.05),
                        Inches(page_h * 0.75),
                        Inches(page_w * 0.9),
                        Inches(page_h * 0.22),
                    )
                    desc_frame = desc_box.text_frame
                    desc_frame.word_wrap = True
                    desc_paragraph = desc_frame.paragraphs[0]
                    desc_run = desc_paragraph.add_run()
                    desc_run.text = desc_text[:300] + ("..." if len(desc_text) > 300 else "")
                    desc_run.font.size = Pt(12)
                    desc_run.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

                if notes_text:
                    slide.notes_slide.notes_text_frame.text = notes_text

        if output_file:
            prs.save(output_file)
            return None

        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes.getvalue()

    @staticmethod
    def create_pdf_from_images(
        image_paths: list[str],
        output_file: str = None,
        aspect_ratio: str = "16:9",
    ) -> Optional[bytes]:
        valid_paths = []
        for path in image_paths:
            if os.path.exists(path):
                valid_paths.append(path)
            else:
                logger.warning("Skipping missing image for PDF export: %s", path)

        if not valid_paths:
            raise ValueError("No valid images found for PDF export")

        try:
            import img2pdf

            page_w, page_h = get_page_size_inches(aspect_ratio)
            layout_fun = img2pdf.get_layout_fun(
                pagesize=(img2pdf.in_to_pt(page_w), img2pdf.in_to_pt(page_h)),
                fit=img2pdf.FitMode.fill,
            )
            pdf_bytes = img2pdf.convert(valid_paths, layout_fun=layout_fun)

            if output_file:
                with open(output_file, "wb") as file_obj:
                    file_obj.write(pdf_bytes)
                return None
            return pdf_bytes
        except Exception as exc:
            logger.warning("img2pdf conversion failed, falling back to Pillow: %s", exc)
            return PptExportService.create_pdf_from_images_pillow(valid_paths, output_file, aspect_ratio)

    @staticmethod
    def create_pdf_from_images_pillow(
        image_paths: list[str],
        output_file: str = None,
        aspect_ratio: str = "16:9",
    ) -> Optional[bytes]:
        from PIL import Image

        page_w, page_h = get_page_size_inches(aspect_ratio)
        width_px = int(page_w * 72)
        height_px = int(page_h * 72)

        images = []
        for path in image_paths:
            img = Image.open(path)
            img = img.convert("RGB")
            img = img.resize((width_px, height_px), Image.Resampling.LANCZOS)
            images.append(img)

        if not images:
            raise ValueError("No valid images found")

        pdf_bytes = io.BytesIO()
        images[0].save(
            pdf_bytes,
            format="PDF",
            save_all=True,
            append_images=images[1:],
            resolution=72.0,
        )
        pdf_bytes.seek(0)

        if output_file:
            with open(output_file, "wb") as file_obj:
                file_obj.write(pdf_bytes.getvalue())
            return None
        return pdf_bytes.getvalue()


_export_service: Optional[PptExportService] = None


def get_ppt_export_service() -> PptExportService:
    global _export_service
    if _export_service is None:
        _export_service = PptExportService()
    return _export_service
