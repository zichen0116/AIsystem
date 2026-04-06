"""
PPT生成模块 - Celery异步任务
"""
import asyncio
import logging
import os
import tempfile
import uuid
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

from celery import Task
from fastapi import UploadFile
from app.celery import celery_app
from app.core.config import get_settings
from app.generators.ppt.banana_service import get_banana_service
from app.generators.ppt.file_service import get_oss_service
from app.generators.ppt.page_utils import (
    build_page_image_prompt,
    extract_page_points,
    get_active_extra_fields_config,
    split_generated_description,
)

logger = logging.getLogger(__name__)

# 本地导出目录（OSS 不可达时的降级存储）
_LOCAL_EXPORTS_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "..", "exports")
_REPO_ROOT = Path(__file__).resolve().parents[4]
_FRONTEND_PUBLIC_DIR = _REPO_ROOT / "teacher-platform" / "public"


def _get_oss_service_safe():
    """获取 OSS 服务，失败时返回 None（降级到本地存储）"""
    try:
        return get_oss_service()
    except Exception as e:
        logger.warning("OSS 服务不可用，降级到本地文件: %s", e)
        return None


def _save_bytes_local(data: bytes, filename: str) -> str:
    """保存字节数据到本地导出目录，返回相对路径标识"""
    os.makedirs(_LOCAL_EXPORTS_DIR, exist_ok=True)
    local_path = os.path.join(_LOCAL_EXPORTS_DIR, filename)
    with open(local_path, "wb") as f:
        f.write(data)
    return local_path


def _upload_or_save_local(oss_svc, data: bytes, oss_key: str, filename: str) -> tuple:
    """
    尝试上传到 OSS；失败时保存到本地。
    返回 (url_or_path, is_local: bool)
    """
    if oss_svc is not None:
        try:
            url = oss_svc.upload_bytes(data, oss_key)
            return url, False
        except Exception as e:
            logger.warning("OSS 上传失败，降级本地存储 key=%s: %s", oss_key, e)
    local_path = _save_bytes_local(data, filename)
    return local_path, True


def _load_project_template_bytes(project_settings: dict | None, oss_svc=None) -> Optional[bytes]:
    settings = project_settings or {}
    template_oss_key = str(settings.get("template_oss_key") or "").strip()
    if template_oss_key and oss_svc is not None:
        try:
            suffix = Path(template_oss_key).suffix or ".png"
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp_path = tmp.name
            try:
                oss_svc.download_file(template_oss_key, tmp_path)
                with open(tmp_path, "rb") as f:
                    return f.read()
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
        except Exception as e:
            logger.warning("加载项目模板图片失败 oss_key=%s: %s", template_oss_key, e)

    template_url = str(settings.get("template_image_url") or "").strip()
    if not template_url:
        return None

    local_bytes = _load_frontend_public_asset_bytes(template_url)
    if local_bytes:
        return local_bytes

    try:
        import urllib.request as _ureq
        with _ureq.urlopen(template_url, timeout=20) as resp:
            return resp.read()
    except Exception as e:
        logger.warning("加载项目模板图片失败 template_url=%s: %s", template_url, e)
        return None


def _load_frontend_public_asset_bytes(asset_url: str) -> Optional[bytes]:
    path = _resolve_frontend_public_asset_path(asset_url)
    if not path:
        return None
    try:
        return path.read_bytes()
    except Exception as e:
        logger.warning("读取前端静态模板失败 asset_url=%s: %s", asset_url, e)
        return None


def _resolve_project_style_prompt(project) -> Optional[str]:
    style = str(getattr(project, "template_style", "") or "").strip()
    if style:
        return style
    settings_style = str((project.settings or {}).get("template_style") or "").strip()
    if settings_style:
        return settings_style
    return project.theme

def _project_settings_with_template_style(project) -> dict:
    settings = dict(project.settings or {})
    style = str(getattr(project, "template_style", "") or "").strip()
    if style:
        settings["template_style"] = style
    return settings



def _resolve_frontend_public_asset_path(asset_url: str) -> Optional[Path]:
    if not asset_url or not asset_url.startswith("/"):
        return None

    parsed = urlparse(asset_url)
    relative_path = parsed.path.lstrip("/")
    if not relative_path:
        return None

    public_root = _FRONTEND_PUBLIC_DIR.resolve()
    candidate = (public_root / relative_path).resolve()
    try:
        candidate.relative_to(public_root)
    except ValueError:
        return None

    if candidate.is_file():
        return candidate
    return None


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

async def _update_task_status(db, task_id_str: str, status: str, progress: int, result: dict = None):
    """更新 PPTTask 记录状态"""
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTTask
    from datetime import datetime, timezone

    res = await db.execute(select(PPTTask).where(PPTTask.task_id == task_id_str))
    task = res.scalar_one_or_none()
    if task:
        task.status = status
        task.progress = progress
        if result is not None:
            task.result = result
        if status in ("COMPLETED", "FAILED"):
            task.completed_at = datetime.now(timezone.utc)
    await db.commit()


def _apply_generated_description(page, generated_text: str, extra_fields_config: list[dict]) -> None:
    parsed = split_generated_description(generated_text, extra_fields_config)
    page.description = parsed["description"] or generated_text

    config = dict(page.config or {})
    stored_extra_fields = dict(config.get("extra_fields") or {})
    managed_keys = [
        str(field.get("key") or "").strip()
        for field in extra_fields_config
        if str(field.get("key") or "").strip() and str(field.get("key") or "").strip() != "notes"
    ]
    for key in managed_keys:
        if key in parsed["extra_fields"]:
            stored_extra_fields[key] = parsed["extra_fields"][key]
        else:
            stored_extra_fields.pop(key, None)
    config["extra_fields"] = stored_extra_fields
    page.config = config

    if any(str(field.get("key") or "").strip() == "notes" for field in extra_fields_config):
        page.notes = parsed["notes"] or None


def _extract_oss_key(url: str) -> Optional[str]:
    """从 OSS URL 或签名 URL 提取 oss_key"""
    try:
        from urllib.parse import urlparse
        parsed = urlparse(url)
        path = parsed.path.lstrip("/")
        if path:
            return path
    except Exception:
        pass
    return None


def _normalize_extension(file_type: Optional[str], filename: Optional[str] = None) -> str:
    """将 file_type/MIME/文件名统一为扩展名（不带点）。"""
    mime_to_ext = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/webp": "webp",
        "application/pdf": "pdf",
        "application/vnd.ms-powerpoint": "ppt",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/msword": "doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "text/plain": "txt",
        "text/markdown": "md",
        "text/csv": "csv",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    }

    value = (file_type or "").strip().lower()
    if value in mime_to_ext:
        return mime_to_ext[value]

    if "/" in value and value in mime_to_ext:
        return mime_to_ext[value]

    if value.startswith("."):
        return value.lstrip(".")

    if value and "/" not in value:
        return value

    if filename:
        ext = os.path.splitext(filename)[1].lower().lstrip(".")
        if ext:
            return ext

    return "bin"


def _load_bytes_from_url_or_path(source: Optional[str]) -> Optional[bytes]:
    """Read bytes from a local file path or remote URL."""
    import urllib.request

    if not source:
        return None

    if os.path.isfile(source):
        try:
            with open(source, "rb") as f:
                return f.read()
        except Exception as e:
            logger.warning("读取本地资源失败 source=%s: %s", source, e)
            return None

    try:
        with urllib.request.urlopen(source, timeout=30) as resp:
            return resp.read()
    except Exception as e:
        logger.warning("下载资源失败 source=%s: %s", source, e)
        return None


async def _read_uploaded_context_images(uploaded_files: list[UploadFile] | None) -> list[bytes]:
    """Read uploaded edit-context images into memory, ignoring non-image files."""
    result: list[bytes] = []
    for upload in uploaded_files or []:
        ext = _normalize_extension(upload.content_type, upload.filename)
        if ext not in ("png", "jpg", "jpeg", "webp", "gif"):
            continue
        data = await upload.read()
        if data:
            result.append(data)
        try:
            await upload.seek(0)
        except Exception:
            pass
    return result


async def _generate_image_with_retry(
    image_provider,
    *,
    prompt: str,
    ref_images: Optional[list[bytes]],
    aspect_ratio: str,
    resolution: str = "2K",
    max_attempts: Optional[int] = None,
    retry_delay: float = 1.0,
) -> bytes:
    """Retry image generation to smooth transient upstream disconnects."""
    settings = get_settings()
    attempts = max_attempts or max(1, int(settings.GENAI_MAX_RETRIES or 0) + 1)
    last_error: Optional[Exception] = None

    for attempt in range(1, attempts + 1):
        try:
            return await image_provider.agenerate_image(
                prompt=prompt,
                ref_images=ref_images or None,
                aspect_ratio=aspect_ratio,
                resolution=resolution,
            )
        except Exception as exc:
            last_error = exc
            if attempt >= attempts:
                raise
            logger.warning("图片生成第 %s/%s 次失败，准备重试: %s", attempt, attempts, exc)
            await asyncio.sleep(max(retry_delay, 0) * attempt)

    if last_error:
        raise last_error
    raise RuntimeError("图片生成失败")


def _download_image_bytes(page, oss_svc) -> Optional[bytes]:
    """同步下载页面图片，返回 bytes；失败返回 None。oss_svc=None 时直接 URL 下载。"""
    if not page.image_url:
        return None
    tmp_path = None
    # 尝试 OSS 下载
    if oss_svc is not None:
        oss_key = _extract_oss_key(page.image_url)
        if oss_key:
            try:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp_path = tmp.name
                oss_svc.download_file(oss_key, tmp_path)
                with open(tmp_path, "rb") as f:
                    return f.read()
            except Exception as e:
                logger.warning("OSS 下载失败 page_id=%s, 降级 URL: %s", page.id, e)
            finally:
                if tmp_path:
                    try:
                        os.unlink(tmp_path)
                    except Exception:
                        pass
    # 直接从 URL / 本地路径下载（fallback）
    return _load_bytes_from_url_or_path(page.image_url)


# ---------------------------------------------------------------------------
# 任务：批量生成描述
# ---------------------------------------------------------------------------

def _download_image_bytes(page, oss_svc) -> Optional[bytes]:
    """Download a page image as bytes, falling back from OSS to URL/local path."""
    if not page.image_url:
        return None

    tmp_path = None
    if oss_svc is not None:
        oss_key = _extract_oss_key(page.image_url)
        if oss_key:
            try:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp_path = tmp.name
                oss_svc.download_file(oss_key, tmp_path)
                with open(tmp_path, "rb") as f:
                    return f.read()
            except Exception as e:
                logger.warning("OSS 涓嬭浇澶辫触 page_id=%s, 闄嶇骇 URL: %s", page.id, e)
            finally:
                if tmp_path:
                    try:
                        os.unlink(tmp_path)
                    except Exception:
                        pass

    return _load_bytes_from_url_or_path(page.image_url)


def _normalize_selection_bbox(
    selection_bbox: Optional[dict],
    *,
    image_width: int,
    image_height: int,
    min_size: int = 8,
) -> Optional[tuple[int, int, int, int]]:
    """Normalize a selection bbox into a clamped (left, top, right, bottom) tuple."""
    import math

    if not isinstance(selection_bbox, dict):
        return None

    try:
        if {"x", "y", "width", "height"}.issubset(selection_bbox):
            x1 = float(selection_bbox["x"])
            y1 = float(selection_bbox["y"])
            x2 = x1 + float(selection_bbox["width"])
            y2 = y1 + float(selection_bbox["height"])
        elif {"x1", "y1", "x2", "y2"}.issubset(selection_bbox):
            x1 = float(selection_bbox["x1"])
            y1 = float(selection_bbox["y1"])
            x2 = float(selection_bbox["x2"])
            y2 = float(selection_bbox["y2"])
        else:
            return None
    except (TypeError, ValueError):
        return None

    left = max(0, min(image_width, math.floor(min(x1, x2))))
    top = max(0, min(image_height, math.floor(min(y1, y2))))
    right = max(0, min(image_width, math.ceil(max(x1, x2))))
    bottom = max(0, min(image_height, math.ceil(max(y1, y2))))

    if right - left < min_size or bottom - top < min_size:
        return None

    return left, top, right, bottom


def _crop_image_by_selection_bbox(
    image_bytes: Optional[bytes],
    selection_bbox: Optional[dict],
) -> Optional[bytes]:
    """Crop the selected region from the current page image for edit reference."""
    import io
    from PIL import Image as PILImage

    if not image_bytes or not selection_bbox:
        return None

    try:
        with PILImage.open(io.BytesIO(image_bytes)) as image:
            crop_box = _normalize_selection_bbox(
                selection_bbox,
                image_width=image.width,
                image_height=image.height,
            )
            if not crop_box:
                return None

            cropped = image.crop(crop_box)
            output = io.BytesIO()
            cropped.save(output, format="PNG")
            return output.getvalue()
    except Exception as e:
        logger.warning("瑁佸壀鍦ㄧ嚎缂栬緫妗嗛€夊尯鍩熷け璐? %s", e)
        return None


@celery_app.task(bind=True, name="banana-slides.generate_descriptions")
def generate_descriptions_task(self: Task, project_id: int, task_id_str: str = None, language: str = "zh"):
    """
    批量生成页面描述。

    Args:
        project_id: 项目ID
        task_id_str: PPTTask.task_id，用于更新进度
        language: 输出语言
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject, PPTPage

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 0)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(
                select(PPTPage)
                .where(PPTPage.project_id == project_id)
                .order_by(PPTPage.page_number)
            )
            pages = list(res.scalars().all())
            total = len(pages)

            if total == 0:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "COMPLETED", 100, {"count": 0})
                return {"status": "completed", "count": 0}

            for p in pages:
                p.is_description_generating = True
            await db.commit()

        banana_svc = get_banana_service()
        theme = _resolve_project_style_prompt(project)
        project_settings = _project_settings_with_template_style(project)
        extra_fields_config = get_active_extra_fields_config(project_settings)
        completed = 0

        for page in pages:
            page_dict = {
                "id": page.id,
                "title": page.title or "",
                "points": extract_page_points(page),
            }

            try:
                description = await banana_svc.generate_description(
                    page_dict,
                    theme=theme,
                    language=language,
                    extra_fields_config=extra_fields_config,
                )
                async with AsyncSessionLocal() as db2:
                    res2 = await db2.execute(select(PPTPage).where(PPTPage.id == page.id))
                    db_page = res2.scalar_one_or_none()
                    if db_page:
                        _apply_generated_description(db_page, description, extra_fields_config)
                        db_page.is_description_generating = False
                    await db2.commit()
            except Exception as e:
                logger.error("描述生成失败 page_id=%s: %s", page.id, e)
                async with AsyncSessionLocal() as db2:
                    res2 = await db2.execute(select(PPTPage).where(PPTPage.id == page.id))
                    db_page = res2.scalar_one_or_none()
                    if db_page:
                        db_page.is_description_generating = False
                    await db2.commit()

            completed += 1
            if task_id_str:
                progress = int(completed / total * 100)
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "PROCESSING", progress)

        if task_id_str:
            async with AsyncSessionLocal() as db4:
                await _update_task_status(db4, task_id_str, "COMPLETED", 100, {"count": completed})
        return {"status": "completed", "project_id": project_id, "count": completed}

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：批量生成图片
# ---------------------------------------------------------------------------

async def generate_images_async(project_id: int, page_ids: list = None, task_id_str: str = None):
    """
    批量生成页面图片（异步核心，可直接 asyncio.create_task 调用，无需 Celery worker）。
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select, update as sa_update
    from app.generators.ppt.banana_models import PPTProject, PPTPage, PPTMaterial, PageImageVersion
    from app.generators.ppt.banana_providers import get_image_provider_singleton

    try:
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 0)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            q = select(PPTPage).where(PPTPage.project_id == project_id)
            if page_ids:
                q = q.where(PPTPage.id.in_(page_ids))
            q = q.order_by(PPTPage.page_number)
            res = await db.execute(q)
            pages = list(res.scalars().all())
            total = len(pages)

            if total == 0:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "COMPLETED", 100, {"count": 0})
                return {"status": "completed", "count": 0}

            for p in pages:
                p.is_image_generating = True
            await db.commit()

        project_settings = _project_settings_with_template_style(project)
        aspect_ratio = project_settings.get("aspect_ratio", "16:9")
        resolution = project_settings.get("image_resolution", "2K")
        oss_svc = _get_oss_service_safe()
        image_provider = get_image_provider_singleton()
        extra_fields_config = get_active_extra_fields_config(project_settings)
        template_ref_image = _load_project_template_bytes(project_settings, oss_svc)
        completed = 0
        success_count = 0
        failure_count = 0
        failed_pages = []

        for page in pages:
            try:
                prompt = build_page_image_prompt(page, extra_fields_config, project_settings)

                # 素材参考图
                ref_images = [template_ref_image] if template_ref_image else []
                if page.material_ids:
                    async with AsyncSessionLocal() as db_mat:
                        mat_res = await db_mat.execute(
                            select(PPTMaterial).where(PPTMaterial.id.in_(page.material_ids))
                        )
                        materials = list(mat_res.scalars().all())
                    for mat in materials:
                        mat_ext = _normalize_extension(mat.file_type, mat.filename)
                        if mat_ext in ("png", "jpg", "jpeg", "webp"):
                            mat_key = _extract_oss_key(mat.url)
                            if mat_key and oss_svc:
                                try:
                                    with tempfile.NamedTemporaryFile(
                                        suffix="." + mat_ext, delete=False
                                    ) as tmp:
                                        tmp_path = tmp.name
                                    oss_svc.download_file(mat_key, tmp_path)
                                    with open(tmp_path, "rb") as f:
                                        ref_images.append(f.read())
                                    os.unlink(tmp_path)
                                except Exception as e:
                                    logger.warning("加载素材图片失败 material_id=%s: %s", mat.id, e)

                img_bytes = await _generate_image_with_retry(
                    image_provider,
                    prompt=prompt,
                    ref_images=ref_images,
                    aspect_ratio=aspect_ratio,
                    resolution=resolution,
                )

                oss_key = "ppt/{}/pages/{}/v{}_{}.png".format(
                    project_id, page.id, page.image_version + 1, uuid.uuid4().hex[:8]
                )
                filename = "page_{}_{}.png".format(page.id, uuid.uuid4().hex[:8])
                image_url, is_local = _upload_or_save_local(oss_svc, img_bytes, oss_key, filename)

                async with AsyncSessionLocal() as db2:
                    res2 = await db2.execute(select(PPTPage).where(PPTPage.id == page.id))
                    db_page = res2.scalar_one_or_none()
                    if db_page:
                        new_version = db_page.image_version + 1
                        db_page.image_url = image_url
                        db_page.image_version = new_version
                        db_page.is_image_generating = False

                        await db2.execute(
                            sa_update(PageImageVersion)
                            .where(PageImageVersion.page_id == page.id)
                            .where(PageImageVersion.is_active.is_(True))
                            .values(is_active=False)
                        )
                        db2.add(PageImageVersion(
                            page_id=page.id,
                            user_id=project.user_id,
                            version=new_version,
                            image_url=image_url,
                            operation="generate",
                            is_active=True,
                        ))
                        await db2.commit()

                        from app.services.redis_service import invalidate_ppt_cover
                        await invalidate_ppt_cover(project_id)
                        success_count += 1

            except Exception as e:
                logger.error("图片生成失败 page_id=%s: %s", page.id, e)
                failure_count += 1
                failed_pages.append({
                    "page_id": page.id,
                    "page_number": page.page_number,
                    "title": page.title,
                    "error": str(e),
                })
                async with AsyncSessionLocal() as db2:
                    res2 = await db2.execute(select(PPTPage).where(PPTPage.id == page.id))
                    db_page = res2.scalar_one_or_none()
                    if db_page:
                        db_page.is_image_generating = False
                    await db2.commit()

            completed += 1
            if task_id_str:
                progress = int(completed / total * 100)
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "PROCESSING", progress)

        """
        prompt = "{0}\n\n编辑要求：{1}".format(base_desc, edit_instruction) if base_desc else edit_instruction

        """
        if task_id_str:
            result = {
                "count": completed,
                "success_count": success_count,
                "failure_count": failure_count,
                "failed_pages": failed_pages,
            }
            async with AsyncSessionLocal() as db4:
                final_status = "FAILED" if success_count == 0 and failure_count > 0 else "COMPLETED"
                await _update_task_status(db4, task_id_str, final_status, 100, result)
        return {
            "status": "failed" if success_count == 0 and failure_count > 0 else "completed",
            "project_id": project_id,
            "count": completed,
            "success_count": success_count,
            "failure_count": failure_count,
            "failed_pages": failed_pages,
        }

    except Exception as e:
        logger.error("generate_images_async 失败 project_id=%s: %s", project_id, e)
        if task_id_str:
            try:
                async with AsyncSessionLocal() as db_err:
                    await _update_task_status(db_err, task_id_str, "FAILED", 0, {"error": str(e)})
            except Exception:
                pass


@celery_app.task(bind=True, name="banana-slides.generate_images")
def generate_images_task(self: Task, project_id: int, page_ids: list = None, task_id_str: str = None):
    """Celery wrapper — delegates to generate_images_async."""
    return asyncio.run(generate_images_async(project_id, page_ids, task_id_str))


# ---------------------------------------------------------------------------
# 共用：下载所有页面图片到临时目录
# ---------------------------------------------------------------------------

def _download_pages_to_tmpdir(pages, oss_svc) -> tuple:
    """下载页面图片到临时目录，返回 (tmp_dir, image_paths)。oss_svc=None 时直接 URL 下载。"""
    import urllib.request
    tmp_dir = tempfile.mkdtemp()
    image_paths = []
    for i, page in enumerate(pages):
        if not page.image_url:
            continue
        local_path = os.path.join(tmp_dir, "page_{:04d}.png".format(i))
        oss_key = _extract_oss_key(page.image_url)
        downloaded = False
        # 优先尝试 OSS 下载
        if oss_svc is not None and oss_key:
            try:
                oss_svc.download_file(oss_key, local_path)
                image_paths.append(local_path)
                downloaded = True
            except Exception as e:
                logger.warning("OSS 下载失败 page_id=%s, 降级 URL: %s", page.id, e)
        # OSS 不可用或失败时，直接从 URL 下载
        if not downloaded:
            try:
                urllib.request.urlretrieve(page.image_url, local_path)
                image_paths.append(local_path)
            except Exception as e:
                logger.warning("下载页面图片失败 page_id=%s: %s", page.id, e)
    return tmp_dir, image_paths


async def _update_project_export(project_id: int, export_url: str):
    """更新项目导出字段"""
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject
    from datetime import datetime, timezone

    async with AsyncSessionLocal() as db:
        res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
        proj = res.scalar_one_or_none()
        if proj:
            proj.exported_file_url = export_url
            proj.exported_at = datetime.now(timezone.utc)
        await db.commit()


# ---------------------------------------------------------------------------
# 任务：导出 PPTX
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.export_pptx")
def export_pptx_task(self: Task, project_id: int, task_id_str: str = None):
    """
    导出PPTX。

    Args:
        project_id: 项目ID
        task_id_str: PPTTask.task_id
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject, PPTPage
    from app.generators.ppt.ppt_export_service import get_ppt_export_service
    import shutil

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 10)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(
                select(PPTPage)
                .where(PPTPage.project_id == project_id)
                .where(PPTPage.image_url.isnot(None))
                .order_by(PPTPage.page_number)
            )
            pages = list(res.scalars().all())

        project_settings = _project_settings_with_template_style(project)
        aspect_ratio = project_settings.get("aspect_ratio", "16:9")
        oss_svc = _get_oss_service_safe()
        export_svc = get_ppt_export_service()

        tmp_dir, image_paths = _download_pages_to_tmpdir(pages, oss_svc)
        try:
            if task_id_str:
                async with AsyncSessionLocal() as db2:
                    await _update_task_status(db2, task_id_str, "PROCESSING", 50)

            pptx_bytes = export_svc.create_pptx_from_images(image_paths, aspect_ratio=aspect_ratio)
            if not pptx_bytes:
                raise ValueError("PPTX生成失败")

            uid = uuid.uuid4().hex[:8]
            oss_key = "ppt/{}/exports/presentation_{}.pptx".format(project_id, uid)
            filename = "presentation_{}_{}.pptx".format(project_id, uid)
            export_url, is_local = _upload_or_save_local(oss_svc, pptx_bytes, oss_key, filename)
            await _update_project_export(project_id, export_url)

            result = {"url": export_url}
            if is_local:
                result["is_local"] = True
                result["warning"] = "OSS不可用，文件已保存到本地，下载时请使用同步接口"
            if task_id_str:
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "COMPLETED", 100, result)
            return {"status": "completed", **result}
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：导出 PDF
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.export_pdf")
def export_pdf_task(self: Task, project_id: int, task_id_str: str = None):
    """
    导出PDF。

    Args:
        project_id: 项目ID
        task_id_str: PPTTask.task_id
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject, PPTPage
    from app.generators.ppt.ppt_export_service import get_ppt_export_service
    import shutil

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 10)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(
                select(PPTPage)
                .where(PPTPage.project_id == project_id)
                .where(PPTPage.image_url.isnot(None))
                .order_by(PPTPage.page_number)
            )
            pages = list(res.scalars().all())

        project_settings = _project_settings_with_template_style(project)
        aspect_ratio = project_settings.get("aspect_ratio", "16:9")
        oss_svc = _get_oss_service_safe()
        export_svc = get_ppt_export_service()

        tmp_dir, image_paths = _download_pages_to_tmpdir(pages, oss_svc)
        try:
            if task_id_str:
                async with AsyncSessionLocal() as db2:
                    await _update_task_status(db2, task_id_str, "PROCESSING", 50)

            pdf_bytes = export_svc.create_pdf_from_images(image_paths, aspect_ratio=aspect_ratio)
            if not pdf_bytes:
                raise ValueError("PDF生成失败")

            uid = uuid.uuid4().hex[:8]
            oss_key = "ppt/{}/exports/presentation_{}.pdf".format(project_id, uid)
            filename = "presentation_{}_{}.pdf".format(project_id, uid)
            export_url, is_local = _upload_or_save_local(oss_svc, pdf_bytes, oss_key, filename)
            await _update_project_export(project_id, export_url)

            result = {"url": export_url}
            if is_local:
                result["is_local"] = True
                result["warning"] = "OSS不可用，文件已保存到本地，下载时请使用同步接口"
            if task_id_str:
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "COMPLETED", 100, result)
            return {"status": "completed", **result}
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：导出图片集 (ZIP)
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.export_images")
def export_images_task(self: Task, project_id: int, task_id_str: str = None):
    """
    将每页图片打包为 ZIP 导出。

    Args:
        project_id: 项目ID
        task_id_str: PPTTask.task_id
    """
    import io as _io
    import zipfile
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject, PPTPage

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 10)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(
                select(PPTPage)
                .where(PPTPage.project_id == project_id)
                .where(PPTPage.image_url.isnot(None))
                .order_by(PPTPage.page_number)
            )
            pages = list(res.scalars().all())

        oss_svc = _get_oss_service_safe()
        zip_buf = _io.BytesIO()
        total = len(pages)
        skipped = 0

        with zipfile.ZipFile(zip_buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(pages):
                img_bytes = _download_image_bytes(page, oss_svc)
                if img_bytes:
                    zf.writestr("slide_{:03d}.png".format(i + 1), img_bytes)
                else:
                    skipped += 1
                    logger.warning("跳过无图片页面 page_id=%s", page.id)

                if task_id_str:
                    progress = int((i + 1) / max(total, 1) * 80) + 10
                    async with AsyncSessionLocal() as db2:
                        await _update_task_status(db2, task_id_str, "PROCESSING", progress)

        zip_bytes = zip_buf.getvalue()
        uid = uuid.uuid4().hex[:8]
        oss_key = "ppt/{}/exports/images_{}.zip".format(project_id, uid)
        filename = "images_{}_{}.zip".format(project_id, uid)
        export_url, is_local = _upload_or_save_local(oss_svc, zip_bytes, oss_key, filename)

        result = {"url": export_url, "total": total, "skipped": skipped}
        if is_local:
            result["is_local"] = True
            result["warning"] = "OSS不可用，文件已保存到本地"
        if task_id_str:
            async with AsyncSessionLocal() as db3:
                await _update_task_status(db3, task_id_str, "COMPLETED", 100, result)
        return {"status": "completed", **result}

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：导出可编辑 PPTX（委托给 export_pptx_task）
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.export_editable_pptx")
def export_editable_pptx_task(self: Task, project_id: int, task_id_str: str = None):
    """
    导出可编辑 PPTX：图片底层 + 文字层（标题/描述/备注）。

    Args:
        project_id: 项目ID
        task_id_str: PPTTask.task_id
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTProject, PPTPage
    from app.generators.ppt.ppt_export_service import get_ppt_export_service
    import shutil

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 10)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(
                select(PPTPage)
                .where(PPTPage.project_id == project_id)
                .where(PPTPage.image_url.isnot(None))
                .order_by(PPTPage.page_number)
            )
            pages = list(res.scalars().all())

        project_settings = _project_settings_with_template_style(project)
        aspect_ratio = project_settings.get("aspect_ratio", "16:9")
        oss_svc = _get_oss_service_safe()
        export_svc = get_ppt_export_service()

        pages_data = [
            {"title": p.title or "", "description": p.description or "", "notes": p.notes or ""}
            for p in pages
        ]

        tmp_dir, image_paths = _download_pages_to_tmpdir(pages, oss_svc)
        try:
            if task_id_str:
                async with AsyncSessionLocal() as db2:
                    await _update_task_status(db2, task_id_str, "PROCESSING", 50)

            pptx_bytes = export_svc.create_pptx_from_images(
                image_paths,
                aspect_ratio=aspect_ratio,
                pages_data=pages_data,
                add_text_layer=True,
            )
            if not pptx_bytes:
                raise ValueError("可编辑PPTX生成失败")

            uid = uuid.uuid4().hex[:8]
            oss_key = "ppt/{}/exports/editable_{}.pptx".format(project_id, uid)
            filename = "editable_{}_{}.pptx".format(project_id, uid)
            export_url, is_local = _upload_or_save_local(oss_svc, pptx_bytes, oss_key, filename)
            await _update_project_export(project_id, export_url)

            result = {"url": export_url}
            if is_local:
                result["is_local"] = True
                result["warning"] = "OSS不可用，文件已保存到本地"
            if task_id_str:
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "COMPLETED", 100, result)
            return {"status": "completed", **result}
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：翻新解析旧PPT/PDF
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.renovation_parse")
def renovation_parse_task(self: Task, project_id: int, file_id: int, task_id_str: str = None):
    """
    解析旧 PPT/PDF 提取大纲。

    Args:
        project_id: 项目ID
        file_id: PPTReferenceFile.id
        task_id_str: PPTTask.task_id
    """
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select
    from app.generators.ppt.banana_models import PPTReferenceFile

    async def _run():
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 0)

            res = await db.execute(select(PPTReferenceFile).where(PPTReferenceFile.id == file_id))
            ref_file = res.scalar_one_or_none()
            if not ref_file:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "File not found"})
                return {"error": "File not found"}

            ref_file.parse_status = "processing"
            await db.commit()

        oss_svc = get_oss_service()
        banana_svc = get_banana_service()
        tmp_path = None

        try:
            ref_ext = _normalize_extension(ref_file.file_type, ref_file.filename)
            with tempfile.NamedTemporaryFile(suffix="." + ref_ext, delete=False) as tmp:
                tmp_path = tmp.name
            oss_key = _extract_oss_key(ref_file.url)
            if oss_key:
                oss_svc.download_file(oss_key, tmp_path)
            else:
                import urllib.request
                urllib.request.urlretrieve(ref_file.url, tmp_path)

            text_content = ""
            if ref_ext == "pdf":
                try:
                    import fitz
                    doc = fitz.open(tmp_path)
                    text_content = "\n".join((page.get_text() or "") for page in doc)
                    doc.close()
                except Exception as e:
                    logger.warning("PDF文本提取失败: %s", e)
            elif ref_ext in ("pptx", "ppt"):
                try:
                    from pptx import Presentation
                    prs = Presentation(tmp_path)
                    lines = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                lines.append(shape.text_frame.text)
                    text_content = "\n".join(lines)
                except Exception as e:
                    logger.warning("PPTX文本提取失败: %s", e)

            if task_id_str:
                async with AsyncSessionLocal() as db2:
                    await _update_task_status(db2, task_id_str, "PROCESSING", 50)

            parsed = await banana_svc.parse_outline_text(
                text_content or ref_file.filename, language="zh"
            )

            async with AsyncSessionLocal() as db3:
                res3 = await db3.execute(
                    select(PPTReferenceFile).where(PPTReferenceFile.id == file_id)
                )
                db_file = res3.scalar_one_or_none()
                if db_file:
                    db_file.parse_status = "completed"
                    db_file.parsed_outline = parsed
                await db3.commit()

            if task_id_str:
                async with AsyncSessionLocal() as db4:
                    await _update_task_status(db4, task_id_str, "COMPLETED", 100, {"parsed": True})
            return {"status": "completed", "file_id": file_id}

        except Exception as e:
            logger.error("翻新解析失败 file_id=%s: %s", file_id, e)
            async with AsyncSessionLocal() as db2:
                res2 = await db2.execute(
                    select(PPTReferenceFile).where(PPTReferenceFile.id == file_id)
                )
                db_file = res2.scalar_one_or_none()
                if db_file:
                    db_file.parse_status = "failed"
                    db_file.parse_error = str(e)
                await db2.commit()
            if task_id_str:
                async with AsyncSessionLocal() as db3:
                    await _update_task_status(db3, task_id_str, "FAILED", 0, {"error": str(e)})
            raise
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

    return asyncio.run(_run())


# ---------------------------------------------------------------------------
# 任务：编辑单页图片
# ---------------------------------------------------------------------------

async def edit_page_image_async(
    project_id: int,
    page_id: int,
    edit_instruction: str,
    task_id_str: str = None,
    context_images: dict = None,
    uploaded_context_images: list[UploadFile] | None = None,
):
    """Async core for single-page image editing, usable without Redis/Celery."""
    from app.core.database import AsyncSessionLocal
    from sqlalchemy import select, update as sa_update
    from app.generators.ppt.banana_models import PPTProject, PPTPage, PageImageVersion, PPTMaterial
    from app.generators.ppt.banana_providers import get_image_provider_singleton

    try:
        async with AsyncSessionLocal() as db:
            if task_id_str:
                await _update_task_status(db, task_id_str, "PROCESSING", 0)

            res = await db.execute(select(PPTProject).where(PPTProject.id == project_id))
            project = res.scalar_one_or_none()
            if not project:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Project not found"})
                return {"error": "Project not found"}

            res = await db.execute(select(PPTPage).where(PPTPage.id == page_id))
            page = res.scalar_one_or_none()
            if not page:
                if task_id_str:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": "Page not found"})
                return {"error": "Page not found"}

            page.is_image_generating = True
            await db.commit()

        project_settings = _project_settings_with_template_style(project)
        aspect_ratio = project_settings.get("aspect_ratio", "16:9")
        resolution = project_settings.get("image_resolution", "2K")
        oss_svc = _get_oss_service_safe()
        image_provider = get_image_provider_singleton()

        ref_images: list[bytes] = []
        page_image_bytes: Optional[bytes] = None
        has_selection_region = False

        if page.image_url:
            page_image_bytes = _download_image_bytes(page, oss_svc)
            if page_image_bytes:
                ref_images.append(page_image_bytes)

        if context_images:
            selection_crop = _crop_image_by_selection_bbox(
                page_image_bytes,
                context_images.get("selection_bbox"),
            )
            if selection_crop:
                ref_images.append(selection_crop)
                has_selection_region = True

            if context_images.get("use_template"):
                template_bytes = _load_project_template_bytes(project.settings or {}, oss_svc)
                if template_bytes:
                    ref_images.append(template_bytes)

            for url in (context_images.get("desc_image_urls") or []):
                img_bytes = _load_bytes_from_url_or_path(url)
                if img_bytes:
                    ref_images.append(img_bytes)

            uploaded_ids = context_images.get("uploaded_image_ids") or []
            if uploaded_ids:
                async with AsyncSessionLocal() as db_mat:
                    mat_res = await db_mat.execute(
                        select(PPTMaterial).where(PPTMaterial.id.in_([
                            int(i) for i in uploaded_ids if str(i).isdigit()
                        ]))
                    )
                    for mat in mat_res.scalars().all():
                        if not mat.url:
                            continue
                        mat_ext = _normalize_extension(mat.file_type, mat.filename)
                        if mat_ext not in ("png", "jpg", "jpeg", "webp"):
                            continue
                        img_bytes = None
                        mat_key = _extract_oss_key(mat.url)
                        if mat_key and oss_svc is not None:
                            tmp_mat_path = None
                            try:
                                with tempfile.NamedTemporaryFile(suffix="." + mat_ext, delete=False) as tmp:
                                    tmp_mat_path = tmp.name
                                oss_svc.download_file(mat_key, tmp_mat_path)
                                with open(tmp_mat_path, "rb") as f:
                                    img_bytes = f.read()
                            except Exception as e:
                                logger.warning("加载上传图片失败 material_id=%s: %s", mat.id, e)
                            finally:
                                if tmp_mat_path:
                                    try:
                                        os.unlink(tmp_mat_path)
                                    except Exception:
                                        pass
                        else:
                            img_bytes = _load_bytes_from_url_or_path(mat.url)
                        if img_bytes:
                            ref_images.append(img_bytes)

        ref_images.extend(await _read_uploaded_context_images(uploaded_context_images))

        base_desc = page.description or page.title or ""
        if has_selection_region:
            edit_instruction = f"{edit_instruction}\n\nPrioritize the user-selected region while keeping the rest of the slide layout and style as consistent as possible."
            """
                f"{edit_instruction}\n\n请优先围绕用户框选的局部区域完成修改，其余区域尽量保持原有版式和风格。"
            )
        prompt = "{0}\n\n编辑要求：{1}".format(base_desc, edit_instruction) if base_desc else edit_instruction

        if task_id_str:
            async with AsyncSessionLocal() as db2:
                await _update_task_status(db2, task_id_str, "PROCESSING", 30)

        """
        prompt = "{0}\n\n编辑要求：{1}".format(base_desc, edit_instruction) if base_desc else edit_instruction

        if task_id_str:
            async with AsyncSessionLocal() as db2:
                await _update_task_status(db2, task_id_str, "PROCESSING", 30)

        img_bytes = await _generate_image_with_retry(
            image_provider,
            prompt=prompt,
            ref_images=ref_images,
            aspect_ratio=aspect_ratio,
            resolution=resolution,
        )

        oss_key = "ppt/{}/pages/{}/edit_{}.png".format(project_id, page_id, uuid.uuid4().hex[:8])
        filename = "page_edit_{}_{}.png".format(page_id, uuid.uuid4().hex[:8])
        image_url, _is_local = _upload_or_save_local(oss_svc, img_bytes, oss_key, filename)

        async with AsyncSessionLocal() as db3:
            res3 = await db3.execute(select(PPTPage).where(PPTPage.id == page_id))
            db_page = res3.scalar_one_or_none()
            if db_page:
                new_version = db_page.image_version + 1
                db_page.image_url = image_url
                db_page.image_version = new_version
                db_page.is_image_generating = False

                await db3.execute(
                    sa_update(PageImageVersion)
                    .where(PageImageVersion.page_id == page_id)
                    .where(PageImageVersion.is_active.is_(True))
                    .values(is_active=False)
                )
                db3.add(PageImageVersion(
                    page_id=page_id,
                    user_id=project.user_id,
                    version=new_version,
                    image_url=image_url,
                    operation="edit",
                    prompt=edit_instruction,
                    is_active=True,
                ))
                await db3.commit()

                from app.services.redis_service import invalidate_ppt_cover
                await invalidate_ppt_cover(project_id)

        if task_id_str:
            async with AsyncSessionLocal() as db4:
                await _update_task_status(db4, task_id_str, "COMPLETED", 100, {"url": image_url})
        return {"status": "completed", "url": image_url}

    except Exception as e:
        logger.error("编辑页面图片失败 page_id=%s: %s", page_id, e)
        async with AsyncSessionLocal() as db_err:
            res = await db_err.execute(select(PPTPage).where(PPTPage.id == page_id))
            db_page = res.scalar_one_or_none()
            if db_page:
                db_page.is_image_generating = False
            if task_id_str:
                await _update_task_status(db_err, task_id_str, "FAILED", 0, {"error": str(e)})
            await db_err.commit()
        return {"status": "failed", "error": str(e)}


@celery_app.task(bind=True, name="banana-slides.edit_page_image")
def edit_page_image_task(
    self: Task,
    project_id: int,
    page_id: int,
    edit_instruction: str,
    task_id_str: str = None,
    context_images: dict = None,
):
    return asyncio.run(
        edit_page_image_async(
            project_id=project_id,
            page_id=page_id,
            edit_instruction=edit_instruction,
            task_id_str=task_id_str,
            context_images=context_images,
        )
    )


# ---------------------------------------------------------------------------
# 任务：AI 生成素材图片
# ---------------------------------------------------------------------------

@celery_app.task(bind=True, name="banana-slides.generate_material")
def generate_material_task(
    self: Task,
    project_id: int,
    user_id: int,
    prompt: str,
    aspect_ratio: str = "1:1",
    task_id_str: str = None,
):
    """
    AI 生成素材图片，结果写入 PPTMaterial 表。

    Args:
        project_id: 项目ID
        user_id: 用户ID
        prompt: 生成提示词
        aspect_ratio: 图片比例
        task_id_str: PPTTask.task_id
    """
    from app.core.database import AsyncSessionLocal
    from app.generators.ppt.banana_models import PPTMaterial
    from app.generators.ppt.banana_providers import get_image_provider_singleton

    async def _run():
        if task_id_str:
            async with AsyncSessionLocal() as db:
                await _update_task_status(db, task_id_str, "PROCESSING", 10)

        try:
            image_provider = get_image_provider_singleton()
            oss_svc = _get_oss_service_safe()

            if task_id_str:
                async with AsyncSessionLocal() as db:
                    await _update_task_status(db, task_id_str, "PROCESSING", 30)

            img_bytes = await image_provider.agenerate_image(
                prompt=prompt,
                aspect_ratio=aspect_ratio,
            )

            uid = uuid.uuid4().hex[:10]
            oss_key = "ppt/{}/materials/ai_{}.png".format(project_id, uid)
            filename = "material_{}_{}.png".format(project_id, uid)
            image_url, _is_local = _upload_or_save_local(oss_svc, img_bytes, oss_key, filename)

            if task_id_str:
                async with AsyncSessionLocal() as db:
                    await _update_task_status(db, task_id_str, "PROCESSING", 80)

            # 写入 PPTMaterial
            async with AsyncSessionLocal() as db:
                material = PPTMaterial(
                    user_id=user_id,
                    project_id=project_id,
                    filename="ai_material_{}.png".format(uuid.uuid4().hex[:6]),
                    oss_path=oss_key,
                    url=image_url,
                    file_type="png",
                    material_type="image",
                )
                db.add(material)
                await db.commit()
                await db.refresh(material)
                material_id = material.id

            if task_id_str:
                async with AsyncSessionLocal() as db:
                    await _update_task_status(
                        db, task_id_str, "COMPLETED", 100,
                        {"url": image_url, "material_id": material_id}
                    )
            return {"status": "completed", "url": image_url, "material_id": material_id}

        except Exception as e:
            logger.error("素材生成失败 project_id=%s: %s", project_id, e)
            if task_id_str:
                async with AsyncSessionLocal() as db:
                    await _update_task_status(db, task_id_str, "FAILED", 0, {"error": str(e)})
            raise

    return asyncio.run(_run())
