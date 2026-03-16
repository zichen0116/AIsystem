"""从 PDF、Word、PPT、文本中提取内容"""
import os
from pathlib import Path
from typing import Optional


def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".txt":
        return _extract_txt(file_path)
    if suffix == ".pdf":
        return _extract_pdf(file_path)
    if suffix in (".doc", ".docx"):
        return _extract_docx(file_path)
    if suffix in (".ppt", ".pptx"):
        return _extract_pptx(file_path)
    return ""


def _extract_txt(file_path: Path) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def _extract_pdf(file_path: Path) -> str:
    """
    PDF 文本提取：
    - 优先使用 PyMuPDF（pymupdf），对中文/编码通常更稳
    - 回退 PyPDF2
    """
    # 1) PyMuPDF (fitz)
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(str(file_path))
        parts = []
        for page in doc:
            text = page.get_text("text") or ""
            text = text.strip()
            if text:
                parts.append(text)
        doc.close()
        content = "\n\n".join(parts).strip()
        if content:
            return content
    except Exception:
        # 继续回退
        pass

    # 2) PyPDF2 fallback
    try:
        import PyPDF2

        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            parts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                text = text.strip()
                if text:
                    parts.append(text)
            return "\n\n".join(parts).strip()
    except Exception:
        return ""


def _extract_docx(file_path: Path) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def _extract_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n\n".join(parts)
    except Exception:
        return ""
